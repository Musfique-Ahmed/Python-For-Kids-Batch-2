from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Content for the slides
slides_content = [
    {
        "title": "E-Commerce & Marketplace AI Adoption Report",
        "content": [
            "Analysis of AI usage in e-commerce platforms",
            "Trends, gaps, and AI adoption roadmap"
        ]
    },
    {
        "title": "Current AI Adoption – Overview",
        "content": [
            "Out of 58 platforms:",
            "- 38 have no AI features",
            "- 20 use some form of AI"
        ]
    },
    {
        "title": "Current AI Adoption – Key Findings",
        "content": [
            "Recommendation Systems: 19 have them, 38 do not",
            "NLP: 14 use it, 44 do not",
            "Personalization: 56 have no visible personalization",
            "Visual Search: 57 do not have it",
            "Fraud Detection: 57 lack it",
            "Inventory Prediction: 56 do not have it",
            "3D Try-out: None",
            "Computer Vision: 56 do not use it",
            "Chatbots: 26 rule-based, 18 AI-powered, 13 none"
        ]
    },
    {
        "title": "MVP 1 – Core AI Features (Quick Wins)",
        "content": [
            "Recommendation System – suggest products based on browsing & purchase history",
            "AI Chatbot with Basic NLP – handle common customer questions",
            "Basic Personalization – show recently viewed items or tailored product categories"
        ]
    },
    {
        "title": "MVP 2 – Better Experience & Efficiency",
        "content": [
            "Advanced NLP – handle complex queries, sentiment analysis, CRM integration",
            "Visual Search – upload images to find similar products",
            "Basic Fraud Detection – flag suspicious transactions",
            "Basic Inventory Prediction – forecast demand to manage stock"
        ]
    },
    {
        "title": "MVP 3 – Advanced AI for Competitive Advantage",
        "content": [
            "3D Try-out – virtual try-on for clothes or furniture",
            "Full Fraud Detection – high-accuracy prevention using multiple data points",
            "Advanced Inventory Optimization – real-time stock management",
            "Computer Vision – auto-categorize products and check image quality"
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "Start small, build gradually",
            "Adopt AI features step-by-step",
            "Gain competitive advantage in the market"
        ]
    }
]

# Create a presentation
prs = Presentation()

# Function to add slides
def add_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.shapes.placeholders[1]

    title_placeholder.text = title

    tf = body_placeholder.text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.name = 'Calibri'
        p.font.color.rgb = RGBColor(0, 0, 0)

# Add all slides
for slide in slides_content:
    add_slide(prs, slide["title"], slide["content"])

# Save the file
output_path = "AI_Adoption_Report.pptx"
prs.save(output_path)
print(f"Presentation saved as {output_path}")
